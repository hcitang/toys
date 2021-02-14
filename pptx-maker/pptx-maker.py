from pptx import Presentation
import csv

def make_presentation(filename, title):
	p = Presentation()
	title_slide = p.slides.add_slide(p.slide_layouts[0])
	title_slide.shapes.title.text = title
	title_slide.placeholders[1].text = "[your names]"

	second_slide = p.slides.add_slide(p.slide_layouts[1])
	second_slide.shapes.title.text = "Mini-Presentation Template"
	tf = second_slide.placeholders[1].text_frame
	tf.text = """Give a summary; do not tell people everything that happened in the paper.

	Slide 1: Title of the paper, the authors, and who your team is
	Slide 2: What is the research question they are asking? Why does it matter?
	Slide 3: What do they do about their research question? What do they build? Or, what kind of user research do they conduct?
	Slide 4: What do they find? What is surprising about this? Or, unexpected?
	Slide 5: What is the major lesson that we can take away from this work?
	Slide 6: What is a discussion question that comes up because of the work?

	Plan to spend about a minute per minute on slides 2-5, and feel free to add if you feel you can make the time. Aim for a 5 minute presentation time; 7 mins max!"""
	p.save(filename)	

with open('data/readings.csv') as csvfile:
	reader = csv.reader(csvfile)
	name = ""
	counter = 1
	for row in reader:
		if (row[0] == name):
			counter = counter + 1
		else:
			name = row[0]
			counter = 1
		filename = name + '-' + str(counter) + '.pptx'
		make_presentation(filename, row[1])





